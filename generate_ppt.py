from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color Palette ──────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x11, 0x17)   # #0D1117 – deep navy
ACCENT_BLUE  = RGBColor(0x23, 0x8C, 0xF5)   # #238CF5 – bright blue
ACCENT_TEAL  = RGBColor(0x00, 0xD4, 0xAA)   # #00D4AA – teal
ACCENT_PURPLE= RGBColor(0x7C, 0x3A, 0xED)   # #7C3AED – violet
ACCENT_ORANGE= RGBColor(0xFF, 0x6B, 0x35)   # #FF6B35 – orange
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xD6, 0xE0)
CARD_BG      = RGBColor(0x16, 0x1B, 0x22)   # #161B22 – card surface
CARD_BORDER  = RGBColor(0x30, 0x36, 0x3D)   # #30363D


def hex2rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ── Helpers ────────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.color.rgb = line_color if line_color else fill_color or CARD_BG
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=Pt(14), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_gradient_accent_bar(slide, left, top, width, height, color: RGBColor):
    """Thin colored accent line."""
    bar = slide.shapes.add_shape(1,
        Inches(left), Inches(top), Inches(width), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_card(slide, left, top, width, height, title, items, icon="", accent=ACCENT_BLUE):
    """Draw a dark card with title + bullet list."""
    # Card background
    card = add_rect(slide, left, top, width, height,
                    fill_color=CARD_BG, line_color=CARD_BORDER, line_width=Pt(0.75))
    # Top accent bar
    add_gradient_accent_bar(slide, left, top, width, 0.04, accent)
    # Icon + title
    title_text = f"{icon}  {title}" if icon else title
    add_text(slide, title_text,
             left + 0.12, top + 0.08, width - 0.24, 0.32,
             font_size=Pt(11), bold=True, color=accent)
    # Bullet items
    y = top + 0.40
    for item in items:
        add_text(slide, f"▸  {item}",
                 left + 0.15, y, width - 0.30, 0.22,
                 font_size=Pt(8.5), color=LIGHT_GRAY)
        y += 0.22


# ── Slide builders ─────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    set_slide_bg(slide, BG_DARK)

    # Decorative circles (simulated with rectangles – circles not needed)
    # Large blurred circle effect via translucent overlapping shapes
    for r, c in [(3.5, RGBColor(0x03, 0x14, 0x28)), (2.5, RGBColor(0x13, 0x09, 0x28))]:
        sh = slide.shapes.add_shape(9,
            Inches(6.5 - r/2), Inches(3.8 - r/2), Inches(r), Inches(r))
        sh.fill.solid()
        sh.fill.fore_color.rgb = c
        sh.line.fill.background()

    # Tag line bar
    add_gradient_accent_bar(slide, 0.5, 1.6, 2.2, 0.05, ACCENT_TEAL)

    # Main heading
    add_text(slide, "AiMitra", 0.5, 1.0, 9.0, 0.9,
             font_size=Pt(48), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, "AI-Powered Assistant Platform", 0.5, 1.75, 9.0, 0.5,
             font_size=Pt(18), bold=False, color=ACCENT_TEAL, align=PP_ALIGN.LEFT)

    # Sub headline
    add_text(slide,
             "One unified workspace for QA Engineers, Developers & Business Analysts\n"
             "powered by 13 AI providers and 30+ intelligent tools.",
             0.5, 2.45, 9.0, 0.9,
             font_size=Pt(12), color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

    # Stats row
    stats = [
        ("13+", "AI Providers"),
        ("30+", "Intelligent Tools"),
        ("4", "Workspaces"),
        ("Multi", "Model Support"),
    ]
    for i, (num, label) in enumerate(stats):
        x = 0.5 + i * 2.45
        add_rect(slide, x, 3.55, 2.2, 0.85, fill_color=CARD_BG, line_color=ACCENT_BLUE, line_width=Pt(0.75))
        add_text(slide, num,  x + 0.1, 3.58, 2.0, 0.42, font_size=Pt(22), bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.1, 4.0,  2.0, 0.35, font_size=Pt(9),  bold=False, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Footer line
    add_gradient_accent_bar(slide, 0, 7.15, 10, 0.02, ACCENT_BLUE)
    add_text(slide, "Demo Presentation  •  June 2026",
             0, 7.18, 10, 0.3, font_size=Pt(9), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_gradient_accent_bar(slide, 0, 0, 10, 0.06, ACCENT_BLUE)

    add_text(slide, "Platform Overview", 0.5, 0.15, 9, 0.5,
             font_size=Pt(28), bold=True, color=WHITE)
    add_gradient_accent_bar(slide, 0.5, 0.65, 1.8, 0.04, ACCENT_TEAL)

    add_text(slide,
             "AiMitra organises all AI capabilities into 4 focused workspaces — "
             "allowing teams to pick the right tool for the right job.",
             0.5, 0.8, 9, 0.5, font_size=Pt(11), color=LIGHT_GRAY)

    workspaces = [
        ("💬", "Chat",        ACCENT_BLUE,
         ["Free-form conversations", "Code & text generation",
          "Web search (Perplexity)", "Image generation (DALL-E 3)"]),
        ("🛠️", "Dev Tools",   ACCENT_TEAL,
         ["DOM Locator Generator", "Test Case & BDD Generator",
          "Code Review / Debug / Convert", "SQL, Regex, Git, DevOps helpers"]),
        ("📋", "JIRA",        ACCENT_PURPLE,
         ["AI Bug Reporter", "Ticket Validator",
          "Comment Generator", "Test Plan Reviewer"]),
        ("📊", "BA Suite",    ACCENT_ORANGE,
         ["User Story & Use Case Gen", "BRD / Gap Analysis",
          "Meeting Summarizer", "Impact & Stakeholder reports"]),
    ]

    for i, (icon, name, color, items) in enumerate(workspaces):
        col = i % 2
        row = i // 2
        x = 0.4 + col * 4.8
        y = 1.55 + row * 2.75
        add_card(slide, x, y, 4.5, 2.55, name, items, icon, color)

    add_gradient_accent_bar(slide, 0, 7.15, 10, 0.02, ACCENT_BLUE)
    add_text(slide, "AiMitra  •  Platform Overview", 0, 7.18, 10, 0.3,
             font_size=Pt(8), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def slide_usecase_chat(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_gradient_accent_bar(slide, 0, 0, 10, 0.06, ACCENT_BLUE)

    add_text(slide, "Use Case 1 — Chat Workspace", 0.5, 0.12, 9, 0.5,
             font_size=Pt(26), bold=True, color=WHITE)
    add_gradient_accent_bar(slide, 0.5, 0.62, 1.5, 0.04, ACCENT_BLUE)

    add_text(slide, "Your AI companion for daily conversations, code help, research & creative tasks.",
             0.5, 0.75, 9, 0.4, font_size=Pt(11), color=LIGHT_GRAY)

    scenarios = [
        ("💬 Text & Code Generation", ACCENT_BLUE,
         ["Ask anything — code, explanations, documentation",
          "Supports all 13 AI providers (GPT-4o, Claude, Gemini, Llama, Grok…)",
          "Switch models mid-session for different perspectives"]),
        ("🌐 Real-Time Web Search", ACCENT_TEAL,
         ["Powered by Perplexity Sonar models",
          "Fetches live documentation, changelogs, and news",
          "Ideal for: 'What's new in Playwright 1.45?'"]),
        ("🎨 Image Generation", ACCENT_ORANGE,
         ["Powered by OpenAI DALL-E 3",
          "Generate UI mockups, icons, or concept art from prompts",
          "Auto-validates model — shows error if wrong model selected"]),
        ("🤖 Model Comparison", ACCENT_PURPLE,
         ["Run the same prompt across multiple models side-by-side",
          "Compare GPT-4o vs Claude vs Gemini output quality",
          "Pick the best model for your specific task"]),
    ]

    for i, (title, color, bullets) in enumerate(scenarios):
        col = i % 2
        row = i // 2
        x = 0.35 + col * 4.75
        y = 1.35 + row * 2.8
        add_card(slide, x, y, 4.4, 2.6, title, bullets, "", color)

    add_gradient_accent_bar(slide, 0, 7.15, 10, 0.02, ACCENT_BLUE)
    add_text(slide, "AiMitra  •  Use Case 1: Chat", 0, 7.18, 10, 0.3,
             font_size=Pt(8), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def slide_usecase_devtools(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_gradient_accent_bar(slide, 0, 0, 10, 0.06, ACCENT_TEAL)

    add_text(slide, "Use Case 2 — Dev Tools Workspace", 0.5, 0.12, 9, 0.5,
             font_size=Pt(26), bold=True, color=WHITE)
    add_gradient_accent_bar(slide, 0.5, 0.62, 1.8, 0.04, ACCENT_TEAL)
    add_text(slide, "15+ specialised developer & QA tools — all AI-accelerated.",
             0.5, 0.75, 9, 0.4, font_size=Pt(11), color=LIGHT_GRAY)

    tools = [
        ("🎯 DOM Locator Generator", ACCENT_TEAL,
         ["Paste HTML → get XPath & CSS selectors instantly",
          "Great for Selenium / Playwright automation scripts"]),
        ("🧪 Test Case Generator", ACCENT_BLUE,
         ["Describe feature → get ready-to-run test cases",
          "Supports Selenium, Playwright, Jest output formats"]),
        ("📋 BDD Generator", ACCENT_PURPLE,
         ["Convert requirements into Gherkin feature files",
          "Given / When / Then scenarios with examples"]),
        ("🔍 Code Review & Explain", ACCENT_ORANGE,
         ["Paste code → get review, security issues, suggestions",
          "Explain complex code in plain English"]),
        ("🐛 Debug & Fix", ACCENT_TEAL,
         ["Paste failing code + error → get root cause + fix",
          "Understands stack traces and runtime errors"]),
        ("🔄 Code Converter", ACCENT_BLUE,
         ["Translate code between languages (Python↔JS, etc.)",
          "Preserve logic while adapting idioms"]),
        ("🗄️ SQL Helper", ACCENT_PURPLE,
         ["Write, optimise, and explain SQL queries",
          "Supports PostgreSQL, MySQL, SQLite dialects"]),
        ("⚙️ DevOps Generator", ACCENT_ORANGE,
         ["Generate Dockerfiles, CI/CD pipelines, k8s YAMLs",
          "GitHub Actions, Jenkins, and more"]),
    ]

    cols = 4
    card_w = 2.28
    card_h = 1.55
    for i, (title, color, bullets) in enumerate(tools):
        col = i % cols
        row = i // cols
        x = 0.2 + col * (card_w + 0.1)
        y = 1.35 + row * (card_h + 0.15)
        add_card(slide, x, y, card_w, card_h, title, bullets, "", color)

    add_gradient_accent_bar(slide, 0, 7.15, 10, 0.02, ACCENT_TEAL)
    add_text(slide, "AiMitra  •  Use Case 2: Dev Tools", 0, 7.18, 10, 0.3,
             font_size=Pt(8), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def slide_usecase_jira(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_gradient_accent_bar(slide, 0, 0, 10, 0.06, ACCENT_PURPLE)

    add_text(slide, "Use Case 3 — JIRA AI Workspace", 0.5, 0.12, 9, 0.5,
             font_size=Pt(26), bold=True, color=WHITE)
    add_gradient_accent_bar(slide, 0.5, 0.62, 1.6, 0.04, ACCENT_PURPLE)
    add_text(slide, "Automate JIRA workflows — create, validate, and enrich tickets using AI.",
             0.5, 0.75, 9, 0.4, font_size=Pt(11), color=LIGHT_GRAY)

    jira_usecases = [
        ("🐛 AI Bug Reporter", ACCENT_PURPLE,
         ["Describe a bug in plain English",
          "AI generates a fully-structured JIRA bug ticket",
          "Auto-fills: Summary, Steps to Reproduce, Expected vs Actual, Severity"]),
        ("✅ Ticket Validator", ACCENT_BLUE,
         ["Paste an existing JIRA ticket",
          "AI checks completeness and clarity",
          "Get actionable improvement suggestions"]),
        ("💬 Comment Generator", ACCENT_TEAL,
         ["Select a ticket → AI drafts professional comments",
          "Update stakeholders with AI-written status notes",
          "Tone-aware: technical, executive, or casual"]),
        ("📝 Test Plan Reviewer", ACCENT_ORANGE,
         ["Upload/paste your test plan",
          "AI reviews coverage gaps and edge cases",
          "Produces a scored review report"]),
        ("📥 Ticket Loader", ACCENT_PURPLE,
         ["Connect to live JIRA instance via API token",
          "Load and analyse tickets directly in AiMitra",
          "Bulk review or summarise sprints"]),
        ("📌 JQL Builder", ACCENT_BLUE,
         ["Describe what you want in English",
          "AI generates optimised JQL query",
          "Supports complex filters, ordering, dates"]),
    ]

    cols = 3
    card_w = 3.0
    card_h = 2.3
    for i, (title, color, bullets) in enumerate(jira_usecases):
        col = i % cols
        row = i // cols
        x = 0.27 + col * (card_w + 0.14)
        y = 1.35 + row * (card_h + 0.15)
        add_card(slide, x, y, card_w, card_h, title, bullets, "", color)

    add_gradient_accent_bar(slide, 0, 7.15, 10, 0.02, ACCENT_PURPLE)
    add_text(slide, "AiMitra  •  Use Case 3: JIRA AI", 0, 7.18, 10, 0.3,
             font_size=Pt(8), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def slide_usecase_ba(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_gradient_accent_bar(slide, 0, 0, 10, 0.06, ACCENT_ORANGE)

    add_text(slide, "Use Case 4 — Business Analysis Suite", 0.5, 0.12, 9, 0.5,
             font_size=Pt(26), bold=True, color=WHITE)
    add_gradient_accent_bar(slide, 0.5, 0.62, 2.0, 0.04, ACCENT_ORANGE)
    add_text(slide, "10 BA tools to accelerate requirements, documentation, and stakeholder communication.",
             0.5, 0.75, 9, 0.4, font_size=Pt(11), color=LIGHT_GRAY)

    ba_tools = [
        ("📖 User Story Generator", ACCENT_ORANGE,
         ["Feature idea → As a / I want / So that stories",
          "With acceptance criteria and story points"]),
        ("✔️ Acceptance Criteria", ACCENT_BLUE,
         ["User story → detailed acceptance criteria",
          "Given/When/Then + DoD checklist"]),
        ("📐 Use Case Generator", ACCENT_TEAL,
         ["Requirement text → formal use case document",
          "Includes actors, flows, and exceptions"]),
        ("📊 Requirements Analyser", ACCENT_PURPLE,
         ["Upload spec doc → AI identifies gaps, ambiguities",
          "Returns structured requirement analysis"]),
        ("🗺️ Process Flow Generator", ACCENT_ORANGE,
         ["Describe a business process",
          "AI generates step-by-step process flow"]),
        ("📄 BRD Generator", ACCENT_BLUE,
         ["Input: feature brief or meeting notes",
          "Output: full Business Requirements Document"]),
        ("🔍 Gap Analysis", ACCENT_TEAL,
         ["Compare AS-IS vs TO-BE states",
          "Identify missing requirements and risks"]),
        ("🎙️ Meeting Summariser", ACCENT_PURPLE,
         ["Paste meeting transcript or notes",
          "AI extracts decisions, action items, owners"]),
        ("📣 Stakeholder Update", ACCENT_ORANGE,
         ["Technical details → non-technical status update",
          "Tailored for exec / business audience"]),
        ("⚡ Impact Analysis", ACCENT_BLUE,
         ["Describe a change → AI maps downstream impacts",
          "Flags risks across systems and teams"]),
    ]

    cols = 5
    card_w = 1.85
    card_h = 1.65
    for i, (title, color, bullets) in enumerate(ba_tools):
        col = i % cols
        row = i // cols
        x = 0.15 + col * (card_w + 0.09)
        y = 1.35 + row * (card_h + 0.18)
        add_card(slide, x, y, card_w, card_h, title, bullets, "", color)

    add_gradient_accent_bar(slide, 0, 7.15, 10, 0.02, ACCENT_ORANGE)
    add_text(slide, "AiMitra  •  Use Case 4: BA Suite", 0, 7.18, 10, 0.3,
             font_size=Pt(8), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def slide_ai_providers(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_gradient_accent_bar(slide, 0, 0, 10, 0.06, ACCENT_BLUE)

    add_text(slide, "13 AI Providers — One Interface", 0.5, 0.12, 9, 0.5,
             font_size=Pt(26), bold=True, color=WHITE)
    add_gradient_accent_bar(slide, 0.5, 0.62, 1.8, 0.04, ACCENT_TEAL)
    add_text(slide,
             "Switch between providers and models in a single click. "
             "AiMitra auto-validates model compatibility with each tool.",
             0.5, 0.75, 9, 0.4, font_size=Pt(11), color=LIGHT_GRAY)

    providers = [
        ("OpenAI",       "GPT-4o, GPT-4 Turbo, o1, o3-mini, DALL-E 3",     ACCENT_BLUE),
        ("Anthropic",    "Claude Opus 4.7, Sonnet 4.6, Haiku 4.5",          ACCENT_TEAL),
        ("Google",       "Gemini 2.0 Flash, Gemini 1.5, Gemma 3",           ACCENT_ORANGE),
        ("Meta / Groq",  "Llama 3.3 70B, Llama 3.1 8B (ultra-fast)",        ACCENT_PURPLE),
        ("Mistral AI",   "Mistral Large, Codestral, Mixtral 8x7B",          ACCENT_BLUE),
        ("DeepSeek",     "DeepSeek V3, DeepSeek R1 (Thinking)",             ACCENT_TEAL),
        ("xAI / Grok",   "Grok 3, Grok 3 Mini, Grok 2",                    ACCENT_ORANGE),
        ("Perplexity",   "Sonar Pro, Sonar — real-time web search",         ACCENT_PURPLE),
        ("Together AI",  "Llama, Qwen 2.5, Mixtral 8x22B",                 ACCENT_BLUE),
        ("Cerebras",     "Llama 3.3 70B Ultra-Fast inference",              ACCENT_TEAL),
        ("OpenRouter",   "200+ models via one API key",                     ACCENT_ORANGE),
        ("Fireworks AI", "Llama, DeepSeek R1, Qwen 2.5",                   ACCENT_PURPLE),
        ("Cohere",       "Command R+, Command R",                           ACCENT_BLUE),
    ]

    cols = 3
    card_w = 2.95
    card_h = 0.72
    for i, (name, models, color) in enumerate(providers):
        col = i % cols
        row = i // cols
        x = 0.25 + col * (card_w + 0.14)
        y = 1.35 + row * (card_h + 0.1)
        add_card(slide, x, y, card_w, card_h, name, [models], "", color)

    add_gradient_accent_bar(slide, 0, 7.15, 10, 0.02, ACCENT_BLUE)
    add_text(slide, "AiMitra  •  AI Providers", 0, 7.18, 10, 0.3,
             font_size=Pt(8), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def slide_demo_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_gradient_accent_bar(slide, 0, 0, 10, 0.06, ACCENT_TEAL)

    add_text(slide, "Demo Walkthrough", 0.5, 0.12, 9, 0.5,
             font_size=Pt(26), bold=True, color=WHITE)
    add_gradient_accent_bar(slide, 0.5, 0.62, 1.3, 0.04, ACCENT_TEAL)
    add_text(slide, "Suggested sequence for a compelling 15-minute demo.",
             0.5, 0.75, 9, 0.4, font_size=Pt(11), color=LIGHT_GRAY)

    steps = [
        ("1", "Chat — Model Selection",
         "Show model switcher. Ask the same question to GPT-4o and Claude Sonnet. Highlight provider logos.", ACCENT_BLUE),
        ("2", "Chat — Web Search",
         "Switch to Perplexity Sonar. Ask: 'Latest Playwright v1.46 features'. Show live results.", ACCENT_TEAL),
        ("3", "Chat — Image Generation",
         "Switch to DALL-E 3. Prompt: 'A futuristic dashboard UI in dark theme'. Show generated image.", ACCENT_ORANGE),
        ("4", "Dev Tools — DOM Locator",
         "Paste a snippet of your app's HTML. Show generated XPath and CSS selectors instantly.", ACCENT_TEAL),
        ("5", "Dev Tools — BDD Generator",
         "Describe a login feature. Generate Gherkin feature file. Copy to clipboard.", ACCENT_PURPLE),
        ("6", "JIRA — AI Bug Reporter",
         "Describe a bug verbally. Show auto-generated JIRA ticket with all fields populated.", ACCENT_PURPLE),
        ("7", "BA Suite — User Story + BRD",
         "Input a feature brief. Generate user story, then escalate to full BRD document.", ACCENT_ORANGE),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        y = 1.35 + i * 0.78
        # Step number badge
        badge = slide.shapes.add_shape(9,
            Inches(0.3), Inches(y + 0.12), Inches(0.42), Inches(0.42))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        add_text(slide, num, 0.3, y + 0.10, 0.42, 0.42,
                 font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Title
        add_text(slide, title, 0.88, y + 0.04, 3.0, 0.30,
                 font_size=Pt(10.5), bold=True, color=color)
        # Description
        add_text(slide, desc, 0.88, y + 0.34, 8.7, 0.38,
                 font_size=Pt(9), color=LIGHT_GRAY)
        # Separator line
        if i < len(steps) - 1:
            add_gradient_accent_bar(slide, 0.5, y + 0.73, 9.0, 0.01, CARD_BORDER)

    add_gradient_accent_bar(slide, 0, 7.15, 10, 0.02, ACCENT_TEAL)
    add_text(slide, "AiMitra  •  Demo Flow", 0, 7.18, 10, 0.3,
             font_size=Pt(8), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def slide_key_differentiators(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_gradient_accent_bar(slide, 0, 0, 10, 0.06, ACCENT_PURPLE)

    add_text(slide, "Why AiMitra?", 0.5, 0.12, 9, 0.5,
             font_size=Pt(26), bold=True, color=WHITE)
    add_gradient_accent_bar(slide, 0.5, 0.62, 1.0, 0.04, ACCENT_PURPLE)
    add_text(slide, "Key differentiators vs using individual AI tools or ChatGPT directly.",
             0.5, 0.75, 9, 0.4, font_size=Pt(11), color=LIGHT_GRAY)

    diffs = [
        ("🔀 Multi-Model Freedom",
         ["No vendor lock-in — switch any model at any time",
          "Compare outputs side-by-side with Model Compare tool",
          "Auto-sync: live model lists fetched from each provider API"],
         ACCENT_BLUE),
        ("🧩 Purpose-Built Tools",
         ["30+ tools crafted for QA / Dev / BA workflows",
          "Each tool has tailored prompts — not generic chat",
          "Structured output: proper JIRA fields, BRD sections, Gherkin"],
         ACCENT_TEAL),
        ("⚡ Smart Guardrails",
         ["Model-capability warnings before you hit errors",
          "DALL-E validation, reasoning-model warnings, small-model alerts",
          "JQL / BDD outputs validated before display"],
         ACCENT_PURPLE),
        ("🏢 Team-Ready",
         ["JIRA integration with live ticket loading",
          "Shareable chat history and exportable outputs",
          "Works on-premise or cloud — bring your own API keys"],
         ACCENT_ORANGE),
    ]

    for i, (title, bullets, color) in enumerate(diffs):
        col = i % 2
        row = i // 2
        x = 0.4 + col * 4.8
        y = 1.45 + row * 2.7
        add_card(slide, x, y, 4.5, 2.5, title, bullets, "", color)

    add_gradient_accent_bar(slide, 0, 7.15, 10, 0.02, ACCENT_PURPLE)
    add_text(slide, "AiMitra  •  Differentiators", 0, 7.18, 10, 0.3,
             font_size=Pt(8), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def slide_thankyou(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Decorative accent
    add_gradient_accent_bar(slide, 0, 0, 10, 0.06, ACCENT_BLUE)
    add_gradient_accent_bar(slide, 0, 7.44, 10, 0.06, ACCENT_TEAL)

    # Central glow
    for r, c in [(4.5, RGBColor(0x03, 0x11, 0x1F)), (3.0, RGBColor(0x00, 0x1A, 0x15))]:
        sh = slide.shapes.add_shape(9,
            Inches(5.0 - r/2), Inches(3.9 - r/2), Inches(r), Inches(r))
        sh.fill.solid(); sh.fill.fore_color.rgb = c
        sh.line.fill.background()

    add_text(slide, "Thank You!", 0, 2.5, 10, 1.0,
             font_size=Pt(52), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_gradient_accent_bar(slide, 3.8, 3.6, 2.4, 0.05, ACCENT_TEAL)
    add_text(slide, "AiMitra — AI for Every Team", 0, 3.75, 10, 0.55,
             font_size=Pt(18), color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

    add_text(slide,
             "Questions? Let's explore any use case live.",
             0, 4.4, 10, 0.45, font_size=Pt(13), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Bottom pills
    pills = ["💬 Chat", "🛠️ Dev Tools", "📋 JIRA AI", "📊 BA Suite"]
    for i, p in enumerate(pills):
        x = 1.0 + i * 2.1
        add_rect(slide, x, 5.3, 1.85, 0.48, fill_color=CARD_BG,
                 line_color=ACCENT_BLUE, line_width=Pt(0.75))
        add_text(slide, p, x, 5.3, 1.85, 0.48,
                 font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── Main ───────────────────────────────────────────────────────────────────────

def build_ppt():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_overview(prs)
    slide_usecase_chat(prs)
    slide_usecase_devtools(prs)
    slide_usecase_jira(prs)
    slide_usecase_ba(prs)
    slide_ai_providers(prs)
    slide_demo_flow(prs)
    slide_key_differentiators(prs)
    slide_thankyou(prs)

    out = "AiMitra_Demo_Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    build_ppt()
